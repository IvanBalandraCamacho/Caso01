"""
RAG Service API
"""

import os
import json
import uuid
import logging
import zlib
import tempfile
from typing import List, Dict, Any, Optional
from pathlib import Path

import redis
import numpy as np
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field, validator
import openai
from openai import OpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer, CrossEncoder

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Config
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379")
EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSION = 1536
CACHE_TTL = 3600  # 1 hour cache

# Init clients
redis_client = redis.from_url(REDIS_URL, decode_responses=True)  # Fix byte handling
openai_client = OpenAI(api_key=OPENAI_API_KEY)

# Init re-ranking model
cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')

# FastAPI app
app = FastAPI(
    title="RAG Service API",
    description="RAG Service for semantic search",
    version="1.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class DocumentMetadata(BaseModel):
    filename: str
    content_type: str
    workspace_id: Optional[str] = None
    user_id: Optional[str] = None

class SearchRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=1000)
    workspace_id: Optional[str] = None
    limit: int = Field(5, ge=1, le=50)  # Max 50 results
    threshold: float = Field(0.7, ge=0.0, le=1.0)
    offset: int = Field(0, ge=0)

    @validator('query')
    def query_not_empty(cls, v):
        if not v.strip():
            raise ValueError('Query cannot be empty')
        return v.strip()

class RAGIngestRequest(BaseModel):
    document_id: str = Field(..., min_length=1)
    workspace_id: str = Field(..., min_length=1)
    content: str = Field(..., min_length=1)
    metadata: Dict[str, Any] = Field(default_factory=dict)
    user_id: Optional[str] = None

    @validator('content')
    def content_not_empty(cls, v):
        if not v.strip():
            raise ValueError('Content cannot be empty')
        return v.strip()

class IngestResponse(BaseModel):
    document_id: str
    chunks_count: int
    status: str

class BatchIngestRequest(BaseModel):
    documents: List[RAGIngestRequest]

class BatchIngestResponse(BaseModel):
    results: List[IngestResponse]
    total_processed: int
    total_chunks: int

class SearchResult(BaseModel):
    document_id: str
    content: str
    score: float
    metadata: Dict[str, Any]

# Utils
def get_embedding(text: str) -> List[float]:
    """Get embedding with cache"""
    import hashlib
    text_hash = hashlib.md5(text.encode()).hexdigest()
    cache_key = f"embedding:{text_hash}"
    
    cached = redis_client.get(cache_key)
    if cached:
        logger.info(f"Embedding cache hit for {text_hash[:8]}")
        return json.loads(cached)
    
    try:
        response = openai_client.embeddings.create(
            input=text,
            model=EMBEDDING_MODEL
        )
        embedding = response.data[0].embedding
        
        redis_client.setex(cache_key, CACHE_TTL, json.dumps(embedding))
        logger.info(f"Embedding generated and cached for {text_hash[:8]}")
        return embedding
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        raise HTTPException(status_code=500, detail="Embedding failed")

def cosine_similarity(a: List[float], b: List[float]) -> float:
    """Cosine similarity"""
    a = np.array(a)
    b = np.array(b)
    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b))

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Smart chunking with langchain"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len
    )
    return text_splitter.split_text(text)

def extract_text_from_file(file: UploadFile) -> str:
    """Extract text from files with cleanup"""
    filename = file.filename.lower()
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
        temp_file.write(file.file.read())
        temp_file_path = temp_file.name
    
    try:
        if filename.endswith('.pdf'):
            import PyPDF2
            with open(temp_file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text

        elif filename.endswith('.docx'):
            import docx
            doc = docx.Document(temp_file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text

        elif filename.endswith(('.xlsx', '.xls')):
            import pandas as pd
            df = pd.read_excel(temp_file_path)
            text = df.to_string(index=False)
            return text

        elif filename.endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(temp_file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            with open(temp_file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content
    
    finally:
        if os.path.exists(temp_file_path):
            os.unlink(temp_file_path)

# Endpoints
@app.post("/ingest_text", response_model=IngestResponse)
async def ingest_text_content(
    request: RAGIngestRequest
):
    """Index text content"""
    try:
        document_id = str(uuid.uuid4())

        text_content = request.content

        if not text_content.strip():
            raise HTTPException(status_code=400, detail="Empty content")

        chunks = chunk_text(text_content)

        chunk_ids = []
        for i, chunk in enumerate(chunks):
            chunk_id = f"{document_id}_chunk_{i}"

            embedding = get_embedding(chunk)

            metadata_dict = {
                **request.metadata,
                "workspace_id": request.workspace_id,
                "chunk_index": i
            }
            if request.user_id is not None:
                metadata_dict["user_id"] = request.user_id
            
            chunk_data = {
                "document_id": document_id,
                "chunk_id": chunk_id,
                "content": zlib.compress(chunk.encode('utf-8')).hex(),
                "embedding": zlib.compress(json.dumps(embedding).encode('utf-8')).hex(),
                "metadata": json.dumps(metadata_dict)
            }

            redis_client.hset(f"chunk:{chunk_id}", mapping=chunk_data)
            chunk_ids.append(chunk_id)

        doc_metadata = {
            "document_id": document_id,
            "workspace_id": request.workspace_id,
            "total_chunks": len(chunks),
            "chunk_ids": json.dumps(chunk_ids),
            **request.metadata
        }
        if request.user_id is not None:
            doc_metadata["user_id"] = request.user_id
        
        redis_client.hset(f"document:{document_id}", mapping=doc_metadata)

        logger.info(f"Indexed doc {document_id} with {len(chunks)} chunks")

        return IngestResponse(
            document_id=document_id,
            chunks_count=len(chunks),
            status="success"
        )

    except Exception as e:
        logger.error(f"Index error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ingest_batch", response_model=BatchIngestResponse)
async def ingest_batch(request: BatchIngestRequest):
    """Batch index documents"""
    try:
        results = []
        total_chunks = 0
        
        for doc_request in request.documents:
            document_id = str(uuid.uuid4())
            text_content = doc_request.content
            
            if not text_content.strip():
                results.append(IngestResponse(
                    document_id=doc_request.document_id,
                    chunks_count=0,
                    status="error",
                    message="Empty content"
                ))
                continue
            
            chunks = chunk_text(text_content)
            chunk_ids = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{document_id}_chunk_{i}"
                embedding = get_embedding(chunk)
                
                metadata_dict = {
                    **doc_request.metadata,
                    "workspace_id": doc_request.workspace_id,
                    "chunk_index": i
                }
                if doc_request.user_id is not None:
                    metadata_dict["user_id"] = doc_request.user_id
                
                chunk_data = {
                    "document_id": document_id,
                    "chunk_id": chunk_id,
                    "content": zlib.compress(chunk.encode('utf-8')).hex(),
                    "embedding": zlib.compress(json.dumps(embedding).encode('utf-8')).hex(),
                    "metadata": json.dumps(metadata_dict)
                }
                
                redis_client.hset(f"chunk:{chunk_id}", mapping=chunk_data)
                chunk_ids.append(chunk_id)
            
            doc_metadata = {
                "document_id": document_id,
                "workspace_id": doc_request.workspace_id,
                "total_chunks": len(chunks),
                "chunk_ids": json.dumps(chunk_ids),
                **doc_request.metadata
            }
            if doc_request.user_id is not None:
                doc_metadata["user_id"] = doc_request.user_id
            
            redis_client.hset(f"document:{document_id}", mapping=doc_metadata)
            
            results.append(IngestResponse(
                document_id=document_id,
                chunks_count=len(chunks),
                status="success"
            ))
            total_chunks += len(chunks)
        
        return BatchIngestResponse(
            results=results,
            total_processed=len(results),
            total_chunks=total_chunks
        )
    
    except Exception as e:
        logger.error(f"Batch error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/search", response_model=List[SearchResult])
async def search_documents(request: SearchRequest):
    """Search documents"""
    try:
        logger.info(f"Search: '{request.query}' | WS: {request.workspace_id} | Thresh: {request.threshold}")
        
        query_embedding = get_embedding(request.query)

        search_pattern = "chunk:*"
        if request.workspace_id:
            relevant_chunks = []

            for doc_key in redis_client.scan_iter("document:*"):
                doc_data = redis_client.hgetall(doc_key)
                workspace_id = doc_data.get("workspace_id")
                
                if workspace_id == request.workspace_id:
                    chunk_ids_str = doc_data.get("chunk_ids")
                    chunk_ids = json.loads(chunk_ids_str)
                    relevant_chunks.extend(chunk_ids)
        else:
            relevant_chunks = [key.decode().split(":", 1)[1] for key in redis_client.scan_iter(search_pattern)]

        results = []
        for chunk_id in relevant_chunks:
            chunk_data = redis_client.hgetall(f"chunk:{chunk_id}")

            if not chunk_data:
                continue

            embedding_str = chunk_data["embedding"]
            content_str = chunk_data["content"]
            document_id_str = chunk_data["document_id"]
            metadata_str = chunk_data["metadata"]

            embedding = json.loads(zlib.decompress(bytes.fromhex(embedding_str)).decode('utf-8'))
            content = zlib.decompress(bytes.fromhex(content_str)).decode('utf-8')
            similarity = cosine_similarity(query_embedding, embedding)

            if similarity >= request.threshold:
                metadata = json.loads(metadata_str)
                results.append(SearchResult(
                    document_id=document_id_str,
                    content=content,
                    score=similarity,
                    metadata=metadata
                ))

        results.sort(key=lambda x: x.score, reverse=True)
        initial_results = results[:request.limit * 2]

        if initial_results:
            query_content_pairs = [[request.query, result.content] for result in initial_results]
            rerank_scores = cross_encoder.predict(query_content_pairs)
            
            for i, result in enumerate(initial_results):
                result.score = (result.score + rerank_scores[i]) / 2
            
            initial_results.sort(key=lambda x: x.score, reverse=True)

        results = initial_results[request.offset:request.offset + request.limit]

        logger.info(f"Search done: {len(results)} results")

        return results

    except Exception as e:
        logger.error(f"Search error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/delete/{document_id}")
async def delete_document(document_id: str):
    """Delete document"""
    try:
        doc_key = f"document:{document_id}"
        doc_data = redis_client.hgetall(doc_key)

        if not doc_data:
            raise HTTPException(status_code=404, detail="Document not found")

        chunk_ids = json.loads(doc_data.get("chunk_ids", "[]"))
        for chunk_id in chunk_ids:
            redis_client.delete(f"chunk:{chunk_id}")

        redis_client.delete(doc_key)

        logger.info(f"Document {document_id} deleted")

        return {"status": "success", "message": f"Document {document_id} deleted"}

    except Exception as e:
        logger.error(f"Delete error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health")
async def health_check():
    """Health check"""
    try:
        redis_client.ping()

        if not OPENAI_API_KEY:
            return {"status": "error", "message": "OPENAI_API_KEY not configured"}

        return {"status": "healthy", "service": "RAG Service"}

    except Exception as e:
        return {"status": "error", "message": str(e)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)